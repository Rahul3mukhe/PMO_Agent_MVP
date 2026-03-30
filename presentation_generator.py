import io
import json
import os
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from llm_providers import generate_text
from schemas import PMOState

_NAVY = RGBColor(0x00, 0x33, 0x66)
_BLUE = RGBColor(0x00, 0x72, 0xbb)
_WHITE = RGBColor(0xff, 0xff, 0xff)
_DARK = RGBColor(0x2d, 0x37, 0x48)

def _build_presentation_prompt(state: PMOState) -> str:
    proj = state.project
    
    # Extract some context from documents if any
    doc_context = ""
    if "risk_registry" in state.docs:
        doc_context += f"Risk Status: {state.docs['risk_registry'].status}. "
        if state.docs["risk_registry"].reasons:
            doc_context += f"Notes: {'; '.join(state.docs['risk_registry'].reasons[:2])}\n"
    
    # Gate context
    gates_str = ""
    for g in state.gates:
        gates_str += f"- {g.gate}: {'PASS' if getattr(g, 'passed', True) else 'FAIL'}\n"

    prompt = f"""
You are a senior PMO Executive assistant creating a mid-project Client Review presentation.
Review the following project details and generate exactly 4-5 PowerPoint slides.

PROJECT DETAILS
Name: {proj.project_name} (ID: {proj.project_id})
Type: {proj.project_type}
Sponsor: {proj.sponsor or 'N/A'}
Financials:
- Estimated Budget: ${proj.estimated_budget or 0:,.2f}
- Actual Consumed: ${proj.actual_budget_consumed or 0:,.2f}
Status: {state.decision or 'IN PROGRESS'}
Gates:
{gates_str}
Context:
{doc_context}
Scope Summary: {proj.scope_summary or 'N/A'}
Risks: {', '.join(proj.known_risks) if proj.known_risks else 'None'}

INSTRUCTIONS
Create professional exactly 4-5 presentation slides covering:
1. Title Slide
2. Executive Summary
3. Financial & Health Status
4. Key Risks & Open Items
5. Next Steps

CRITICAL INSTRUCTIONS:
- DO NOT HALLUCINATE or invent dates, budgets, or facts.
- Use a strictly formal, executive corporate tone suitable for senior reviewers.

Output MUST be strictly valid JSON matching this schema exactly. DO NOT output ANY other text or markdown blocks outside the JSON:
{{
  "slides": [
    {{
      "title": "Title of Slide",
      "bullets": [
        "First major point",
        "Second major point",
        "Third point"
      ]
    }}
  ]
}}
"""
    return prompt.strip()

def _add_nttdata_logo(slide, prs):
    # Attempt to load nttdata_logo.png
    logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "nttdata_logo.png")
    if not os.path.exists(logo_path):
        logo_path = "nttdata_logo.png"
    
    if os.path.exists(logo_path):
        # Add to top right corner
        left = prs.slide_width - Inches(1.8)
        top = Inches(0.2)
        height = Inches(0.4)
        try:
            slide.shapes.add_picture(logo_path, left, top, height=height)
        except Exception as e:
            print(f"Failed to add logo: {e}")

def generate_client_pptx(state: PMOState) -> bytes:
    prompt = _build_presentation_prompt(state)
    
    # Call LLM
    try:
        response_text = generate_text(
            provider=state.provider,
            model=state.model,
            prompt=prompt,
            api_key=state.audit.get("api_key"),
            temperature=0.0,
            max_tokens=2048,
            standards=state.standards,
            project=state.project,
            doc_type="presentation"
        )
        
        # Clean JSON markdown blocks if any
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
            
        data = json.loads(response_text)
        slides_data = data.get("slides", [])
    except Exception as e:
        print(f"PPTX LLM generation failed: {e}")
        # Fallback slides if LLM completely fails
        slides_data = [
            {"title": "Project Review", "bullets": ["Error generating presentation content.", str(e)]}
        ]

    # Build Presentation
    prs = Presentation()
    
    # Find Title layout (0) and Bullet layout (1)
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    
    for i, slide_info in enumerate(slides_data):
        is_title_slide = (i == 0)
        slide = prs.slides.add_slide(title_layout if is_title_slide else bullet_layout)
        
        # Add background to Title Slide
        if is_title_slide:
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = _NAVY
        else:
            # Add a subtle corporate blue divider line below titles
            from pptx.enum.shapes import MSO_SHAPE
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(1.2), prs.slide_width - Inches(1.0), Pt(2)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = _BLUE
            line.line.fill.background()
            
        _add_nttdata_logo(slide, prs)
        
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get("title", f"Slide {i+1}")
        
        # Style Title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = 'Calibri'
            if is_title_slide:
                paragraph.font.color.rgb = _WHITE
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.font.color.rgb = _NAVY
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
        
        # Content
        if is_title_slide:
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = "\\n".join(slide_info.get("bullets", []))
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.color.rgb = _WHITE
                    paragraph.font.size = Pt(20)
                    paragraph.alignment = PP_ALIGN.CENTER
        else:
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                # slightly move text body down to clear the title line
                body_shape.top = Inches(1.5)
                tf = body_shape.text_frame
                tf.clear()
                for bullet in slide_info.get("bullets", []):
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.name = 'Calibri'
                    p.font.size = Pt(18)
                    p.font.color.rgb = _DARK
                    p.space_after = Pt(14)

    # Add a Footer to all slides
    for slide in prs.slides:
        txBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), prs.slide_width - Inches(1), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{state.standards.get('org', {}).get('name', 'PMO Agent')} - Client Confidential"
        p.font.size = Pt(10)
        p.font.color.rgb = _BLUE
        p.alignment = PP_ALIGN.CENTER
        
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
